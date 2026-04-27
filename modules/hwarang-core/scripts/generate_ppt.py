"""Hwarang AI 홍보 PPT 생성 스크립트."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 브랜드 컬러
PRIMARY = RGBColor(0x63, 0x66, 0xF1)      # Indigo
PRIMARY_DARK = RGBColor(0x43, 0x38, 0xCA)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)       # Purple
DARK = RGBColor(0x0F, 0x17, 0x2A)         # Slate 900
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            line_info = {"text": line_info}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_info.get("text", "")
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.alignment = line_info.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space", 6))
    return txBox


# ════════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK)

# 그라디언트 느낌 장식
add_shape(slide, Inches(-2), Inches(-2), Inches(6), Inches(6), RGBColor(0x1E, 0x1B, 0x4B))
add_shape(slide, Inches(9), Inches(4), Inches(6), Inches(6), RGBColor(0x1E, 0x1B, 0x4B))

# 로고
logo_shape = add_shape(slide, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5), PRIMARY)
add_text(slide, Inches(5.9), Inches(1.55), Inches(1.5), Inches(1.5), "H", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.3), Inches(9.3), Inches(1), "Hwarang AI", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(0.6), "화랑 AI — 한국형 AI 어시스턴트", size=24, color=GRAY, align=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(3), Inches(5.2), Inches(7.3), Inches(1.5), [
    {"text": "코딩 · 법률 · 세무 특화 AI | 토큰 기반 경제 | GPU 공유 네트워크", "size": 14, "color": GRAY, "align": PP_ALIGN.CENTER},
    {"text": "", "size": 8},
    {"text": "(주)퍼시스모어  |  hwarang.ai  |  2026", "size": 12, "color": RGBColor(0x47, 0x55, 0x69), "align": PP_ALIGN.CENTER},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 2: 비전 & 미션
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "VISION & MISSION", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1), "한국에서 쓸 만한 AI를,\n모두와 함께 만들어갑니다.", size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(1),
    "해외 AI에 의존하지 않는, 한국어에 특화된 AI를 만듭니다.\n혼자가 아닌 커뮤니티와 함께 — GPU를 나누고, 토큰을 공유하며, 함께 키웁니다.",
    size=16, color=GRAY)

# 3개 카드
cards = [
    {"icon": "🎯", "title": "한국어 특화", "desc": "코딩, 법률, 세무 3대 도메인\n한국어 최적화 학습"},
    {"icon": "🤝", "title": "함께 만드는 AI", "desc": "GPU 공유 네트워크로\n누구나 기여하고 보상받는 구조"},
    {"icon": "💰", "title": "합리적 가격", "desc": "토큰 기반 종량제\nFree 플랜으로 무료 시작"},
]

for i, card in enumerate(cards):
    x = Inches(0.8 + i * 4)
    card_shape = add_shape(slide, x, Inches(4.0), Inches(3.5), Inches(2.8), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x + Inches(0.3), Inches(4.2), Inches(3), Inches(0.6), card["icon"], size=32)
    add_text(slide, x + Inches(0.3), Inches(4.8), Inches(3), Inches(0.5), card["title"], size=18, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(5.4), Inches(3), Inches(1), card["desc"], size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 3: 제품 라인업
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "PRODUCT LINEUP", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "다양한 플랫폼에서 만나는 화랑 AI", size=32, color=WHITE, bold=True)

products = [
    {"icon": "🌐", "name": "웹 채팅", "desc": "hwarang.ai에서 바로 사용\n스트리밍 대화, 코드 하이라이팅", "status": "✅ 완료", "color": GREEN},
    {"icon": "💻", "name": "VS Code 확장", "desc": "에디터 안에서 AI 코딩\nClaude Code 스타일 UX", "status": "🔧 개발중", "color": CYAN},
    {"icon": "⌨️", "name": "CLI 에이전트", "desc": "터미널 기반 AI 에이전트\nReAct 루프 + 도구 시스템", "status": "🔧 개발중", "color": CYAN},
    {"icon": "🔌", "name": "REST API", "desc": "OpenAI 호환 API\n/v1/chat/completions", "status": "✅ 완료", "color": GREEN},
    {"icon": "🖥️", "name": "데스크톱 앱", "desc": "Windows / macOS 전용 앱\nTauri 기반 네이티브", "status": "📋 계획", "color": GRAY},
    {"icon": "📱", "name": "모바일 앱", "desc": "iOS / Android\nReact Native", "status": "📋 계획", "color": GRAY},
]

for i, p in enumerate(products):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4)
    y = Inches(2.2 + row * 2.5)
    card_shape = add_shape(slide, x, y, Inches(3.5), Inches(2.0), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5), p["icon"], size=24)
    add_text(slide, x + Inches(1.0), y + Inches(0.2), Inches(2.2), Inches(0.4), p["name"], size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(2.9), Inches(0.8), p["desc"], size=11, color=GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(1.6), Inches(2), Inches(0.3), p["status"], size=10, color=p["color"])


# ════════════════════════════════════════════════════════════════
# SLIDE 4: 기술 스택 & 아키텍처
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "TECHNOLOGY", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "검증된 기술 스택", size=32, color=WHITE, bold=True)

# 아키텍처 다이어그램 (텍스트 기반)
arch_box = add_shape(slide, Inches(0.8), Inches(2.0), Inches(7), Inches(5), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(1.2), Inches(2.2), Inches(6), Inches(4.5), [
    {"text": "시스템 아키텍처", "size": 16, "color": WHITE, "bold": True, "space": 12},
    {"text": "┌─────────────────────────────────────────┐", "size": 10, "color": GRAY},
    {"text": "│  사용자: 웹 / VS Code / CLI / API / 앱   │", "size": 10, "color": CYAN},
    {"text": "├─────────────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│  Next.js 15  ←→  FastAPI (OpenAI 호환)  │", "size": 10, "color": PRIMARY},
    {"text": "├─────────────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│  vLLM (Qwen 32B)  +  LoRA 어댑터        │", "size": 10, "color": GREEN},
    {"text": "│  코딩 / 법률 / 세무 도메인 특화          │", "size": 10, "color": GREEN},
    {"text": "├─────────────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│  PostgreSQL + Redis + Prisma             │", "size": 10, "color": ACCENT},
    {"text": "├─────────────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│  Hwarang Grid (GPU 공유 네트워크)        │", "size": 10, "color": RGBColor(0xF5, 0x9E, 0x0B)},
    {"text": "└─────────────────────────────────────────┘", "size": 10, "color": GRAY},
])

# 기술 스택 카드
techs = [
    {"cat": "AI / ML", "items": "Qwen2.5-32B\nQLoRA (bitsandbytes)\nvLLM 서빙\nPyTorch"},
    {"cat": "Backend", "items": "FastAPI\nPostgreSQL\nRedis\nPrisma ORM"},
    {"cat": "Frontend", "items": "Next.js 15\nReact 19\nTailwindCSS 4\nTypeScript"},
    {"cat": "Infra", "items": "RTX 5090 32GB\nNginx + SSL\nCloudflare\nDocker"},
]

for i, tech in enumerate(techs):
    x = Inches(8.3)
    y = Inches(2.0 + i * 1.3)
    card_shape = add_shape(slide, x, y, Inches(4.2), Inches(1.1), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(1.5), Inches(0.3), tech["cat"], size=11, color=PRIMARY, bold=True)
    add_text(slide, x + Inches(1.8), y + Inches(0.1), Inches(2.2), Inches(0.9), tech["items"], size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: 토큰 이코노미
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "TOKEN ECONOMY", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "토큰 기반 경제 시스템", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(0.5), "사용한 만큼만 지불. GPU를 나누면 토큰을 벌 수 있습니다.", size=15, color=GRAY)

plans = [
    {"name": "Free", "price": "무료", "tokens": "10,000", "features": "채팅 기본\n하루 3,000 제한", "color": GRAY},
    {"name": "Starter", "price": "₩9,900/월", "tokens": "100,000", "features": "7B 모델\nAPI 1개\n이메일 지원", "color": CYAN},
    {"name": "Pro", "price": "₩29,900/월", "tokens": "500,000", "features": "32B 모델\nAPI 5개\n우선 지원", "color": PRIMARY},
    {"name": "Business", "price": "₩99,000/월", "tokens": "2,000,000", "features": "전체 모델\nAPI 무제한\n전담 지원", "color": ACCENT},
]

for i, plan in enumerate(plans):
    x = Inches(0.8 + i * 3.1)
    card_shape = add_shape(slide, x, Inches(2.8), Inches(2.8), Inches(4.2), RGBColor(0x1E, 0x29, 0x3B))

    # 플랜 헤더
    header = add_shape(slide, x, Inches(2.8), Inches(2.8), Inches(0.8), plan["color"])
    add_text(slide, x, Inches(2.85), Inches(2.8), Inches(0.35), plan["name"], size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.2), Inches(2.8), Inches(0.35), plan["price"], size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # 토큰
    add_text(slide, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(0.3), "월 토큰", size=10, color=GRAY)
    add_text(slide, x + Inches(0.3), Inches(4.2), Inches(2.2), Inches(0.5), plan["tokens"], size=24, color=WHITE, bold=True)

    # 기능
    add_text(slide, x + Inches(0.3), Inches(5.0), Inches(2.2), Inches(1.5), plan["features"], size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 6: Hwarang Grid (GPU 공유)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "HWARANG GRID", size=12, color=RGBColor(0xF5, 0x9E, 0x0B), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "GPU 공유로 함께 성장하는 네트워크", size=32, color=WHITE, bold=True)

add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), [
    {"text": "내 PC의 유휴 GPU를 공유하고,", "size": 18, "color": WHITE, "bold": True, "space": 4},
    {"text": "토큰을 보상으로 받으세요.", "size": 18, "color": WHITE, "bold": True, "space": 16},
    {"text": ""},
    {"text": "💻  데스크톱 트레이 앱으로 간편 참여", "size": 14, "color": GRAY, "space": 8},
    {"text": "🔒  작업 데이터 암호화 — 내 PC에 저장 안 됨", "size": 14, "color": GRAY, "space": 8},
    {"text": "⚡  처리량 기반 공정 보상 시스템", "size": 14, "color": GRAY, "space": 8},
    {"text": "🎁  연속 참여 보너스 (스트릭 보상)", "size": 14, "color": GRAY, "space": 8},
    {"text": "📊  커뮤니티 페이지에서 실시간 현황 확인", "size": 14, "color": GRAY, "space": 8},
])

# 보상 예시
reward_box = add_shape(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(4.2), [
    {"text": "GPU별 예상 보상", "size": 16, "color": WHITE, "bold": True, "space": 16},
    {"text": "RTX 4060  →  시간당 ~50 토큰", "size": 13, "color": GRAY, "space": 8},
    {"text": "RTX 4070  →  시간당 ~80 토큰", "size": 13, "color": GRAY, "space": 8},
    {"text": "RTX 4080  →  시간당 ~120 토큰", "size": 13, "color": GRAY, "space": 8},
    {"text": "RTX 4090  →  시간당 ~200 토큰", "size": 13, "color": GRAY, "space": 8},
    {"text": "RTX 5090  →  시간당 ~350 토큰", "size": 13, "color": CYAN, "space": 16},
    {"text": ""},
    {"text": "벌어든 토큰으로:", "size": 12, "color": WHITE, "bold": True, "space": 8},
    {"text": "• AI 서비스 무료 이용", "size": 12, "color": GRAY, "space": 4},
    {"text": "• 플랜 업그레이드", "size": 12, "color": GRAY, "space": 4},
    {"text": "• 토큰 거래 (예정)", "size": 12, "color": GRAY, "space": 4},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 7: 도메인 특화
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "DOMAIN EXPERTISE", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "3대 특화 도메인", size=32, color=WHITE, bold=True)

domains = [
    {
        "icon": "⌨️", "name": "코딩", "color": PRIMARY,
        "desc": "• Python, JavaScript, Go 등 전 언어 지원\n• 한국어 주석/설명 자동 생성\n• 코드 리뷰, 리팩토링, 디버깅\n• VS Code 확장으로 에디터 내 사용",
    },
    {
        "icon": "⚖️", "name": "법률", "color": ACCENT,
        "desc": "• 한국 민법, 형법, 노동법 학습\n• 판례 검색 및 요약\n• 계약서 검토 어시스턴트\n• 내용증명 작성 도우미",
    },
    {
        "icon": "🧾", "name": "세무", "color": GREEN,
        "desc": "• 종합소득세, 부가세, 양도세\n• 연말정산 공제 항목 안내\n• 개인사업자 절세 전략\n• 세무 일정 관리",
    },
]

for i, d in enumerate(domains):
    x = Inches(0.8 + i * 4)
    card_shape = add_shape(slide, x, Inches(2.2), Inches(3.6), Inches(4.8), RGBColor(0x1E, 0x29, 0x3B))

    header = add_shape(slide, x, Inches(2.2), Inches(3.6), Inches(1.2), d["color"])
    add_text(slide, x, Inches(2.3), Inches(3.6), Inches(0.5), d["icon"], size=32, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.9), Inches(3.6), Inches(0.4), d["name"], size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), Inches(3.7), Inches(3), Inches(3), d["desc"], size=12, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8: 화랑 고유 정렬 프레임워크 (개요)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HWARANG ALIGNMENT FRAMEWORK", size=12, color=CYAN, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "세계 최초 한국 특화 AI 정렬 기법", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "Claude의 Constitutional AI를 뛰어넘는 10가지 독자 기법. 한국인을 위한, 한국인에 의한, 한국 AI.",
    size=14, color=GRAY)

# 10가지 기법 그리드 (5 x 2)
alignment_all = [
    {"code": "KCAI", "name": "Korean Constitutional AI", "desc": "한국형 헌법 자기비평", "color": PRIMARY},
    {"code": "GRPO", "name": "Grid-based RLHF", "desc": "네트워크 피드백 학습", "color": CYAN},
    {"code": "TACS", "name": "Trust-Aware Safety", "desc": "안전 체인 + 기관 연계", "color": GREEN},
    {"code": "HRAG", "name": "Hwarang RAG", "desc": "한국 공식 DB 실시간 검색", "color": RGBColor(0xF5, 0x9E, 0x0B)},
    {"code": "NWNC", "name": "Noonchi Conversation", "desc": "눈치/정서 기반 대화", "color": RGBColor(0xEC, 0x48, 0x99)},
    {"code": "CoRD", "name": "Collaborative Debate", "desc": "다중 모델 토론", "color": RGBColor(0xF9, 0x73, 0x16)},
    {"code": "VCoT", "name": "Verified CoT", "desc": "각 단계 검증 추론", "color": RGBColor(0x14, 0xB8, 0xA6)},
    {"code": "TADM", "name": "Temporal Decision", "desc": "시점 인식 + 일정", "color": RGBColor(0xA8, 0x55, 0xF7)},
    {"code": "MMRM", "name": "Multi-Memory", "desc": "계층적 개인화 메모리", "color": RGBColor(0x06, 0xB6, 0xD4)},
    {"code": "LCRG", "name": "Citation Guard", "desc": "실시간 인용 검증", "color": RGBColor(0xEF, 0x44, 0x44)},
]

for i, m in enumerate(alignment_all):
    col = i % 5
    row = i // 5
    x = Inches(0.5 + col * 2.55)
    y = Inches(2.7 + row * 2.0)

    card = add_shape(slide, x, y, Inches(2.4), Inches(1.8), RGBColor(0x1E, 0x29, 0x3B))

    # 코드 배지
    code_bg = add_shape(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.9), Inches(0.35), m["color"])
    add_text(slide, x + Inches(0.2), y + Inches(0.18), Inches(0.9), Inches(0.3), m["code"],
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 이름
    add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.1), Inches(0.4), m["name"],
             size=11, color=WHITE, bold=True)

    # 설명
    add_text(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.1), Inches(0.7), m["desc"],
             size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8b: 정렬 기법 상세 (3개 핵심)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "ALIGNMENT SPOTLIGHT", size=12, color=CYAN, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "세계 최초 독자 기법 3선", size=32, color=WHITE, bold=True)

spotlight = [
    {
        "code": "GRPO", "name": "Grid-based RLHF",
        "subtitle": "네트워크 기반 지속 학습",
        "color": CYAN,
        "desc": "• Grid 참여자 = 피드백 제공자\n• 피드백 품질에 따라 토큰 보상\n• 👎 받은 답변 → 자동 DPO 데이터\n• 주기적 LoRA 재학습 & 배포\n\n▶ 한국 실사용자 피드백으로 진화하는\n   세계 유일의 AI 학습 루프",
    },
    {
        "code": "NWNC", "name": "Noonchi Conversation",
        "subtitle": "\"눈치\" 기반 대화 이해",
        "color": RGBColor(0xEC, 0x48, 0x99),
        "desc": "• \"괜찮긴 한데...\" → 불만 파악\n• \"ㅠㅠ\" → 공감 응답 자동\n• 존댓말/반말 상황 자동 전환\n• 체면/배려 문화 반영\n\n▶ 한국어 정서를 이해하는\n   유일한 AI (서양 AI 불가능)",
    },
    {
        "code": "CoRD", "name": "Collaborative Debate",
        "subtitle": "다중 모델 토론 시스템",
        "color": RGBColor(0xF9, 0x73, 0x16),
        "desc": "• 여러 모델이 같은 질문 답변\n• Round 2: 서로의 답 비판/수정\n• Round 3: 심판 모델이 합의 선정\n• 환각(hallucination) 70% 감소\n\n▶ Grid 분산 추론 네트워크에서만\n   가능한 고유 기법",
    },
]

for i, m in enumerate(spotlight):
    x = Inches(0.5 + i * 4.2)
    card = add_shape(slide, x, Inches(2.2), Inches(3.9), Inches(4.8), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.3), Inches(2.4), Inches(1.2), Inches(0.5), m["color"])
    add_text(slide, x + Inches(0.3), Inches(2.45), Inches(1.2), Inches(0.4), m["code"],
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), Inches(3.1), Inches(3.5), Inches(0.4), m["name"],
             size=15, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.5), Inches(3.5), Inches(0.4), m["subtitle"],
             size=11, color=m["color"])

    add_text(slide, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(2.8), m["desc"],
             size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8c: 정렬 파이프라인 흐름도
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "ALIGNMENT PIPELINE", size=12, color=CYAN, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "모든 요청에 자동 적용되는 10단계 파이프라인", size=28, color=WHITE, bold=True)

# 파이프라인 스텝
pipeline_steps = [
    {"step": "1", "name": "질문 수신", "stage": "INPUT", "color": GRAY, "desc": "사용자 메시지"},
    {"step": "2", "name": "NWNC", "stage": "분석", "color": RGBColor(0xEC, 0x48, 0x99), "desc": "눈치/정서 분석"},
    {"step": "3", "name": "TACS", "stage": "분석", "color": GREEN, "desc": "도메인 감지"},
    {"step": "4", "name": "TADM", "stage": "분석", "color": RGBColor(0xA8, 0x55, 0xF7), "desc": "시점 해석"},
    {"step": "5", "name": "MMRM", "stage": "조회", "color": RGBColor(0x06, 0xB6, 0xD4), "desc": "유저 메모리"},
    {"step": "6", "name": "HRAG", "stage": "조회", "color": RGBColor(0xF5, 0x9E, 0x0B), "desc": "한국 공식 DB"},
    {"step": "7", "name": "KCAI", "stage": "프롬프트", "color": PRIMARY, "desc": "헌법 주입"},
    {"step": "8", "name": "LLM", "stage": "생성", "color": CYAN, "desc": "DeepSeek V3 등"},
    {"step": "9", "name": "VCoT + LCRG", "stage": "검증", "color": RGBColor(0x14, 0xB8, 0xA6), "desc": "추론+인용 검증"},
    {"step": "10", "name": "TACS + GRPO", "stage": "출력", "color": RGBColor(0xEF, 0x44, 0x44), "desc": "면책+피드백"},
]

for i, step in enumerate(pipeline_steps):
    col = i % 5
    row = i // 5
    x = Inches(0.5 + col * 2.55)
    y = Inches(2.3 + row * 2.3)

    # 단계 번호 원
    num_bg = add_shape(slide, x + Inches(0.2), y, Inches(0.5), Inches(0.5), step["color"])
    add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(0.5), Inches(0.4), step["step"],
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 카드
    card = add_shape(slide, x, y + Inches(0.7), Inches(2.4), Inches(1.4), RGBColor(0x1E, 0x29, 0x3B))

    add_text(slide, x + Inches(0.2), y + Inches(0.75), Inches(2.1), Inches(0.3), step["stage"],
             size=9, color=step["color"], bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.1), Inches(0.4), step["name"],
             size=13, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.5), Inches(2.1), Inches(0.5), step["desc"],
             size=10, color=GRAY)

    # 화살표 (같은 행 내)
    if col < 4 and i != 9:
        arrow = add_shape(slide, x + Inches(2.4), y + Inches(1.3), Inches(0.15), Pt(2), step["color"])


# ════════════════════════════════════════════════════════════════
# SLIDE 8d: 최적화 프레임워크 (Optimization Stack)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HWARANG OPTIMIZATION STACK", size=12, color=RGBColor(0xF9, 0x73, 0x16), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "추론 속도 + 학습 효율 극대화", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "정렬 10기법에 이어 추론/학습 최적화 7기법 추가. 속도 3배, 비용 50%↓, 학습 효율 40%↑",
    size=14, color=GRAY)

# 7가지 최적화 기법
optimization_stack = [
    {"code": "HPC", "name": "Prompt Cache", "desc": "비용 50% 절감\n속도 3배", "color": CYAN, "category": "추론"},
    {"code": "HAT", "name": "Agentic Tools", "desc": "계산기/법령검색/\n코드실행 도구", "color": RGBColor(0xF9, 0x73, 0x16), "category": "추론"},
    {"code": "HSD", "name": "Speculative Decode", "desc": "Draft 모델로\n속도 2-3배", "color": GREEN, "category": "추론"},
    {"code": "HLC", "name": "Long Context", "desc": "32K → 128K\nYaRN scaling", "color": RGBColor(0xA8, 0x55, 0xF7), "category": "추론"},
    {"code": "HMM", "name": "Model Merging", "desc": "Qwen+EXAONE\n학습없이 합성", "color": PRIMARY, "category": "학습"},
    {"code": "HCL", "name": "Curriculum Learn", "desc": "쉬움→어려움\n단계적 학습", "color": RGBColor(0xEC, 0x48, 0x99), "category": "학습"},
    {"code": "HSR", "name": "Self-Rewarding", "desc": "AI가 스스로\nDPO 데이터 생성", "color": RGBColor(0xEF, 0x44, 0x44), "category": "학습"},
]

# 상단: 추론 최적화 (4개)
inference_opts = [o for o in optimization_stack if o["category"] == "추론"]
add_text(slide, Inches(0.8), Inches(2.5), Inches(5), Inches(0.4), "⚡ 추론 최적화", size=14, color=CYAN, bold=True)

for i, m in enumerate(inference_opts):
    x = Inches(0.8 + i * 3.05)
    card = add_shape(slide, x, Inches(3.0), Inches(2.85), Inches(1.8), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.2), Inches(3.15), Inches(0.9), Inches(0.35), m["color"])
    add_text(slide, x + Inches(0.2), Inches(3.18), Inches(0.9), Inches(0.3), m["code"],
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), Inches(3.6), Inches(2.6), Inches(0.4), m["name"],
             size=12, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(4.05), Inches(2.6), Inches(0.7), m["desc"],
             size=10, color=GRAY)

# 하단: 학습 최적화 (3개)
training_opts = [o for o in optimization_stack if o["category"] == "학습"]
add_text(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.4), "🧠 학습 최적화", size=14, color=RGBColor(0xEC, 0x48, 0x99), bold=True)

for i, m in enumerate(training_opts):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(5.5), Inches(3.9), Inches(1.8), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.2), Inches(5.65), Inches(0.9), Inches(0.35), m["color"])
    add_text(slide, x + Inches(0.2), Inches(5.68), Inches(0.9), Inches(0.3), m["code"],
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), Inches(6.1), Inches(3.6), Inches(0.4), m["name"],
             size=12, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(6.55), Inches(3.6), Inches(0.7), m["desc"],
             size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8f: 고급 추론 기법 (2025~2026 최신)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "ADVANCED TECHNIQUES", size=12, color=RGBColor(0xA8, 0x55, 0xF7), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "2025~2026 최신 학계 기법 도입", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "o1/o3, DeepSeek-R1, Google Chain-of-Agents, MIT TTT 등 최첨단 논문 기법을 화랑에 통합",
    size=14, color=GRAY)

advanced_techniques = [
    {"code": "TTT", "name": "Test-Time Training", "subtitle": "MIT 2025", "desc": "추론 시점 일시 학습\n정확도 30%↑"},
    {"code": "CoA", "name": "Chain-of-Agents", "subtitle": "Google 2025", "desc": "긴 문서 릴레이 처리\n정확도 ↑ 비용 ↓"},
    {"code": "RSL", "name": "Reasoning Scaling", "subtitle": "o1/o3 style", "desc": "깊이 있는 사고\n복잡 추론 특화"},
    {"code": "ACT", "name": "Adaptive Compute", "subtitle": "난이도별 조절", "desc": "쉬움=빠르게\n어려움=깊게"},
    {"code": "MTP", "name": "Multi-Token Pred", "subtitle": "Meta 2024", "desc": "한번에 N토큰\n속도 4배"},
    {"code": "QSR", "name": "Quiet-STaR", "subtitle": "숨은 사고", "desc": "내부 독백으로\n사고 품질↑"},
    {"code": "dGRPO", "name": "DeepSeek GRPO", "subtitle": "R1 핵심", "desc": "그룹 상대 최적화\n메모리 50%↓"},
    {"code": "CDPO", "name": "Constitutional DPO", "subtitle": "Anthropic", "desc": "헌법 기반\n자동 쌍 생성"},
    {"code": "PRM", "name": "Process Reward", "subtitle": "단계별 보상", "desc": "각 단계 평가\n수학/코딩 최고"},
    {"code": "DDP", "name": "Distillation", "subtitle": "증류 학습", "desc": "V3 → Qwen\n지식 전이"},
]

for i, m in enumerate(advanced_techniques):
    col = i % 5
    row = i // 5
    x = Inches(0.5 + col * 2.55)
    y = Inches(2.5 + row * 2.3)

    card = add_shape(slide, x, y, Inches(2.4), Inches(2.0), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.9), Inches(0.35),
                        RGBColor(0xA8, 0x55, 0xF7))
    add_text(slide, x + Inches(0.2), y + Inches(0.18), Inches(0.9), Inches(0.3), m["code"],
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.1), Inches(0.35), m["name"],
             size=11, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.95), Inches(2.1), Inches(0.3), m["subtitle"],
             size=9, color=RGBColor(0xA8, 0x55, 0xF7))
    add_text(slide, x + Inches(0.2), y + Inches(1.3), Inches(2.1), Inches(0.6), m["desc"],
             size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8g: 화랑 독자 혁신 (세계 최초)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HWARANG INNOVATION", size=12, color=RGBColor(0xFB, 0xBF, 0x24), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "세계 최초 독자 기법 6선", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "논문/특허 가능한 수준. Claude/GPT도 아직 못 구현한 한국 AI 독자 혁신.",
    size=14, color=GRAY)

innovations = [
    {
        "code": "HNTL", "name": "Neural Topic Lock",
        "desc": "도메인별 전용 뉴런 경로\n(Mixture-of-LoRAs)\n→ 속도+정확도 동시 향상",
        "color": RGBColor(0xFB, 0xBF, 0x24),
    },
    {
        "code": "HCE", "name": "Community Evolution",
        "desc": "Grid 피드백으로 AI 진화\n(유전 알고리즘 스타일)\n→ 한국 취향 실시간 적응",
        "color": RGBColor(0xEC, 0x48, 0x99),
    },
    {
        "code": "HRL", "name": "Reality Lock",
        "desc": "실시간 팩트 체크\n(문장별 자동 검증)\n→ 환각 원천 차단",
        "color": RGBColor(0xEF, 0x44, 0x44),
    },
    {
        "code": "HCP", "name": "Confidence Pruning",
        "desc": "3개 모델 앙상블\n(logprob 기반 선택)\n→ 환각 90% 감소",
        "color": RGBColor(0x06, 0xB6, 0xD4),
    },
    {
        "code": "HML", "name": "Memory Ladder",
        "desc": "계단식 RAG\n(개념→법령→판례→예외)\n→ 법률 조사 자동화",
        "color": PRIMARY,
    },
    {
        "code": "HQL", "name": "Quantum Learning",
        "desc": "여러 LoRA 동시 학습\n(가중 합성 추론)\n→ 다중 페르소나",
        "color": RGBColor(0xA8, 0x55, 0xF7),
    },
]

for i, m in enumerate(innovations):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.5 + row * 2.3)

    card = add_shape(slide, x, y, Inches(3.9), Inches(2.0), RGBColor(0x1E, 0x29, 0x3B))

    # 별 표시 (세계 최초)
    add_text(slide, x + Inches(3.4), y + Inches(0.15), Inches(0.4), Inches(0.3), "⭐",
             size=14, align=PP_ALIGN.CENTER)

    code_bg = add_shape(slide, x + Inches(0.3), y + Inches(0.2), Inches(1.2), Inches(0.45), m["color"])
    add_text(slide, x + Inches(0.3), y + Inches(0.23), Inches(1.2), Inches(0.35), m["code"],
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.4), Inches(0.35), m["name"],
             size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.4), Inches(0.8), m["desc"],
             size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8e: 경쟁 비교 (화랑 vs Claude vs GPT)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "COMPETITIVE ADVANTAGE", size=12, color=CYAN, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "경쟁사 대비 화랑 AI의 우위", size=32, color=WHITE, bold=True)

# 비교 테이블
compare_headers = ["기능", "Claude", "GPT-4", "화랑 AI"]
compare_rows = [
    ("Constitutional AI", "✅", "❌", "✅ KCAI (한국형)"),
    ("RLHF 학습", "직원", "직원", "✅ GRPO (커뮤니티)"),
    ("한국 공식 DB 연동", "❌", "❌", "✅ HRAG + LCRG"),
    ("한국어 눈치 이해", "약함", "약함", "✅ NWNC"),
    ("실시간 팩트 체크", "❌", "❌", "✅ HRL (세계 최초)"),
    ("다중 모델 앙상블", "❌", "❌", "✅ CoRD + HCP"),
    ("도메인 뉴런 라우팅", "❌", "❌", "✅ HNTL (세계 최초)"),
    ("Test-Time Training", "❌", "❌", "✅ TTT"),
    ("Reasoning (o1 스타일)", "✅", "✅ o1/o3", "✅ RSL + Quiet-STaR"),
    ("Process Reward Model", "✅", "✅", "✅ PRM + VCoT"),
    ("긴 문서 처리", "200K", "128K", "✅ 128K~1M + CoA"),
    ("커뮤니티 진화", "❌", "❌", "✅ HCE (세계 최초)"),
]

# 헤더
header_y = Inches(2.3)
col_widths = [Inches(3.5), Inches(2.5), Inches(2.5), Inches(3.5)]
x = Inches(0.8)
for i, header in enumerate(compare_headers):
    header_bg = add_shape(slide, x, header_y, col_widths[i], Inches(0.5), CYAN if i == 3 else RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x, header_y + Inches(0.1), col_widths[i], Inches(0.3), header,
             size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x += col_widths[i]

# 데이터 행
for row_idx, row in enumerate(compare_rows):
    y = Inches(2.8 + row_idx * 0.42)
    x = Inches(0.8)
    row_bg_color = RGBColor(0x14, 0x1C, 0x2A) if row_idx % 2 == 0 else RGBColor(0x1E, 0x29, 0x3B)

    for i, cell in enumerate(row):
        cell_bg = add_shape(slide, x, y, col_widths[i], Inches(0.4), row_bg_color)
        color = CYAN if i == 3 else WHITE if i == 0 else GRAY
        size = 10
        add_text(slide, x + Inches(0.15), y + Inches(0.07), col_widths[i] - Inches(0.2), Inches(0.3), cell,
                 size=size, color=color, bold=(i == 0 or i == 3))
        x += col_widths[i]


# ════════════════════════════════════════════════════════════════
# SLIDE 9: 현재 진행 상황
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "PROGRESS", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "현재 개발 진행 상황", size=32, color=WHITE, bold=True)

milestones = [
    {"phase": "Phase 1", "title": "기반 구축", "status": "완료", "color": GREEN,
     "items": "모노레포 구조 설계\n공유 스키마 정의\nDB 스키마 (Prisma)"},
    {"phase": "Phase 2", "title": "학습 파이프라인", "status": "진행중", "color": RGBColor(0xF5, 0x9E, 0x0B),
     "items": "Qwen2.5-32B QLoRA 학습 중\nDeepSeek-V3 준비\nSFT/DPO 데이터 수집"},
    {"phase": "Phase 3", "title": "웹 + API", "status": "완료", "color": GREEN,
     "items": "Next.js 웹 채팅\n소셜/이메일 로그인\n관리자 패널\nvLLM OpenAI 호환 API"},
    {"phase": "Phase 4", "title": "정렬 프레임워크", "status": "완료", "color": GREEN,
     "items": "KCAI (한국형 헌법)\nGRPO (Grid RLHF)\nTACS (안전 체인)\n피드백 수집 API"},
    {"phase": "Phase 5", "title": "확장", "status": "계획", "color": GRAY,
     "items": "VS Code 확장팩\nCLI 에이전트\nHwarang Grid 런칭\n데스크톱/모바일"},
]

for i, m in enumerate(milestones):
    x = Inches(0.5 + i * 2.5)
    # 상태 원
    circle = add_shape(slide, x + Inches(0.8), Inches(2.3), Inches(0.5), Inches(0.5), m["color"])
    add_text(slide, x + Inches(0.8), Inches(2.33), Inches(0.5), Inches(0.5),
             "✓" if m["status"] == "완료" else "◎" if m["status"] == "진행중" else "○",
             size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # 연결선 (마지막 제외)
    if i < len(milestones) - 1:
        line = add_shape(slide, x + Inches(1.3), Inches(2.47), Inches(2.0), Pt(2), RGBColor(0x33, 0x41, 0x55))

    card = add_shape(slide, x, Inches(3.0), Inches(2.2), Inches(3.8), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x + Inches(0.2), Inches(3.1), Inches(1.8), Inches(0.3), m["phase"], size=10, color=m["color"], bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.4), Inches(1.8), Inches(0.4), m["title"], size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.9), Inches(1.8), Inches(0.3), m["status"], size=11, color=m["color"])
    add_text(slide, x + Inches(0.2), Inches(4.3), Inches(1.8), Inches(2.2), m["items"], size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 9: 하드웨어
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "INFRASTRUCTURE", size=12, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "서버 인프라", size=32, color=WHITE, bold=True)

hw_box = add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(4.2), [
    {"text": "마스터 서버 (hwarang.ai)", "size": 16, "color": WHITE, "bold": True, "space": 16},
    {"text": "CPU    AMD Ryzen 9 9950X3D (16코어/32스레드)", "size": 12, "color": GRAY, "space": 6},
    {"text": "GPU    NVIDIA RTX 5090 32GB GDDR7", "size": 12, "color": CYAN, "space": 6},
    {"text": "RAM    128GB DDR5-6400", "size": 12, "color": GRAY, "space": 6},
    {"text": "SSD    NVMe 4TB × 3 (학습/모델/서비스 분리)", "size": 12, "color": GRAY, "space": 16},
    {"text": ""},
    {"text": "소프트웨어 스택", "size": 14, "color": WHITE, "bold": True, "space": 12},
    {"text": "OS     Ubuntu 24.04 LTS", "size": 12, "color": GRAY, "space": 6},
    {"text": "CUDA   12.8 + cuDNN 9.x", "size": 12, "color": GRAY, "space": 6},
    {"text": "Proxy  Nginx + Cloudflare + Let's Encrypt", "size": 12, "color": GRAY, "space": 6},
])

scale_box = add_shape(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(4.2), [
    {"text": "확장 계획", "size": 16, "color": WHITE, "bold": True, "space": 16},
    {"text": "마스터/서브 분산 아키텍처", "size": 13, "color": GRAY, "space": 8},
    {"text": "gRPC 기반 워커 통신", "size": 13, "color": GRAY, "space": 8},
    {"text": "자동 로드밸런싱 & 장애 복구", "size": 13, "color": GRAY, "space": 8},
    {"text": "Hwarang Grid 연동", "size": 13, "color": GRAY, "space": 16},
    {"text": ""},
    {"text": "처리 능력 (RTX 5090 + vLLM)", "size": 14, "color": WHITE, "bold": True, "space": 12},
    {"text": "32B INT4 단일 요청:  40~60 토큰/초", "size": 12, "color": CYAN, "space": 6},
    {"text": "배치 처리량:          200~300 토큰/초", "size": 12, "color": CYAN, "space": 6},
    {"text": "동시 요청 (PagedAttn): 32~64건", "size": 12, "color": GRAY, "space": 6},
    {"text": "응답 시간 (짧은답):   0.5~2초", "size": 12, "color": GRAY, "space": 6},
    {"text": "응답 시간 (긴답):      5~15초", "size": 12, "color": GRAY, "space": 6},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 10: HFL - 특허 출원 기술 (핵심 슬라이드)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

# 특허 배지
patent_bg = add_shape(slide, Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.5), RGBColor(0xF5, 0x9E, 0x0B))
add_text(slide, Inches(10.5), Inches(0.33), Inches(2.2), Inches(0.4), "📋 특허 출원 예정",
         size=11, color=DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "PATENT PENDING", size=12, color=RGBColor(0xF5, 0x9E, 0x0B), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1.0),
    "HFL: Federated LoRA Training for\nCommunity-Driven LLM Improvement\nover Low-Bandwidth Networks",
    size=28, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(0.5),
    "세계 최초: GPU 공유 네트워크가 AI를 분산 학습시키는 구조. 일반 인터넷만으로 가능.",
    size=14, color=CYAN)

# 좌측: 문제 → 해결
problem_box = add_shape(slide, Inches(0.5), Inches(3.2), Inches(6), Inches(4.0), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(0.9), Inches(3.3), Inches(5.3), Inches(3.8), [
    {"text": "기존 분산 학습의 한계", "size": 14, "color": RGBColor(0xEF, 0x44, 0x44), "bold": True, "space": 12},
    {"text": "❌ 매 Step마다 전체 Gradient 동기화 필요", "size": 11, "color": GRAY, "space": 6},
    {"text": "❌ InfiniBand 400Gbps 전용 네트워크 필수", "size": 11, "color": GRAY, "space": 6},
    {"text": "❌ 장비 비용 수억원 (H100 클러스터)", "size": 11, "color": GRAY, "space": 6},
    {"text": "❌ 개인/중소기업 접근 불가능", "size": 11, "color": GRAY, "space": 16},
    {"text": ""},
    {"text": "HFL Adaptive의 해결", "size": 14, "color": GREEN, "bold": True, "space": 12},
    {"text": "✅ 5단계 압축: 50MB → 1.8MB (96.5% 감소)", "size": 11, "color": WHITE, "space": 6},
    {"text": "✅ 네트워크 속도 자동 감지 → LoRA 랭크 적응", "size": 11, "color": WHITE, "space": 6},
    {"text": "✅ 1Mbps 모바일에서도 16초 이내 전송", "size": 11, "color": WHITE, "space": 6},
    {"text": "✅ 데이터 프라이버시 (양자화+희소화 비복원)", "size": 11, "color": WHITE, "space": 6},
])

# 우측: 구조도
arch_box = add_shape(slide, Inches(6.8), Inches(3.2), Inches(6), Inches(4.0), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(7.2), Inches(3.3), Inches(5.3), Inches(3.8), [
    {"text": "HFL 동작 원리", "size": 14, "color": WHITE, "bold": True, "space": 12},
    {"text": "┌────────────────────────────────┐", "size": 10, "color": GRAY},
    {"text": "│  마스터 서버 (hwarang.ai)      │", "size": 10, "color": CYAN},
    {"text": "│  LoRA 합성 (TIES/DARE)          │", "size": 10, "color": CYAN},
    {"text": "├────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│  ↑ 1.8MB   ↑ 1.2MB   ↑ 0.5MB  │", "size": 10, "color": GREEN},
    {"text": "├────────────────────────────────┤", "size": 10, "color": GRAY},
    {"text": "│ Worker 1  Worker 2  Worker 3   │", "size": 10, "color": RGBColor(0xF5, 0x9E, 0x0B)},
    {"text": "│ 코딩 LoRA 법률 LoRA 세무 LoRA   │", "size": 10, "color": RGBColor(0xF5, 0x9E, 0x0B)},
    {"text": "│ (RTX 4060) (RTX 4090) (RTX 5090)│", "size": 10, "color": GRAY},
    {"text": "└────────────────────────────────┘", "size": 10, "color": GRAY},
    {"text": "", "space": 8},
    {"text": "통신: 일반 인터넷  |  라운드: 10분/회", "size": 10, "color": WHITE},
    {"text": "참여자: Hwarang Grid 유저 → 토큰 보상", "size": 10, "color": GREEN},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 10b: HFL 에이전트 혁신 (특허 추가 청구항)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HFL AGENT INNOVATION", size=12, color=RGBColor(0xF5, 0x9E, 0x0B), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "HFL 에이전트 3대 혁신 기능", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "특허 청구항 11~13: 에이전트가 더 똑똑해지면 학습 효율이 극대화됩니다.",
    size=14, color=GRAY)

agent_innovations = [
    {
        "code": "ADG", "name": "Auto Data Generation",
        "subtitle": "청구항 11: 자체 데이터 생성",
        "color": RGBColor(0x10, 0xB9, 0x81),
        "desc": "• 마스터에서 키워드만 수신 (1KB)\n• 에이전트가 질문 + 답변 자동 생성\n• Self-Rewarding으로 품질 평가\n• Smart Selection: loss 높은 것만 학습\n\n→ 데이터 배포 완전 불필요!\n→ 프라이버시 완벽 보장",
    },
    {
        "code": "PVA", "name": "Peer Validation",
        "subtitle": "청구항 12: 분산 동료 검증",
        "color": RGBColor(0xEF, 0x44, 0x44),
        "desc": "• 워커 A가 LoRA 제출\n• 랜덤 워커 B, C가 검증\n• baseline loss vs LoRA loss 비교\n• 다수결로 합격/불합격 판정\n\n→ 악성 LoRA 자동 필터링\n→ Poisoning attack 방어",
    },
    {
        "code": "NAS", "name": "Network-Aware Scheduling",
        "subtitle": "청구항 13: 능력 기반 할당",
        "color": RGBColor(0x06, 0xB6, 0xD4),
        "desc": "• GPU 벤치마크 + 네트워크 측정\n• RTX 5090 → 800 step\n• RTX 4060 → 200 step\n• 빠른 워커 완료 → 추가 배치\n\n→ 이종 환경 최적화\n→ Straggler 자동 처리",
    },
]

for i, m in enumerate(agent_innovations):
    x = Inches(0.5 + i * 4.2)
    card = add_shape(slide, x, Inches(2.5), Inches(3.9), Inches(4.7), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.3), Inches(2.7), Inches(1.2), Inches(0.45), m["color"])
    add_text(slide, x + Inches(0.3), Inches(2.73), Inches(1.2), Inches(0.35), m["code"],
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(3.5), Inches(0.35), m["name"],
             size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.65), Inches(3.5), Inches(0.35), m["subtitle"],
             size=10, color=m["color"])

    add_text(slide, x + Inches(0.3), Inches(4.1), Inches(3.5), Inches(2.9), m["desc"],
             size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10c: HFL 인프라 혁신 (장애재할당, 분산서빙, 자가학습)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HFL INFRASTRUCTURE", size=12, color=RGBColor(0xF5, 0x9E, 0x0B), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "멈추지 않는 AI: 장애 복구 + 분산 서빙 + 자가 학습", size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "특허 청구항 14~16: 데이터 손실 0, GPU 부족 해결, 24/7 자동 진화",
    size=14, color=GRAY)

infra_items = [
    {
        "code": "FO", "name": "Task Failover",
        "subtitle": "청구항 14: 장애 재할당",
        "color": RGBColor(0xEF, 0x44, 0x44),
        "desc": "• Heartbeat 3회 실패 → 장애 판정\n• 다른 워커에 즉시 재할당\n• 체크포인트로 이어서 학습\n• 선착순 채택 (중복 허용)\n• 악성 워커 자동 블랙리스트\n\n→ 데이터 손실 0% 보장",
    },
    {
        "code": "DS", "name": "Distributed Serving",
        "subtitle": "청구항 15: 분산 추론 서빙",
        "color": RGBColor(0x06, 0xB6, 0xD4),
        "desc": "• LoRA (2MB)를 Grid 에이전트에 배포\n• 지역/부하/지연 기반 라우팅\n• 한국/일본/미국 에이전트 분산\n• 에이전트 다운 → 자동 페일오버\n• 토큰 보상 (추론 처리 대가)\n\n→ GPU 부족해도 무한 확장",
    },
    {
        "code": "CL", "name": "Continuous Learning",
        "subtitle": "청구항 16: 24/7 자가 학습",
        "color": RGBColor(0x10, 0xB9, 0x81),
        "desc": "• 매일: 뉴스/법령 자동 수집\n• 매주: 유저 피드백 DPO 학습\n• 매월: 전체 벤치마크 + 재학습\n• 이벤트: 법개정 → 즉시 학습\n• 검증 → 자동 배포 → 반복\n\n→ 사람 개입 없이 24/7 진화",
    },
]

for i, m in enumerate(infra_items):
    x = Inches(0.5 + i * 4.2)
    card = add_shape(slide, x, Inches(2.5), Inches(3.9), Inches(4.7), RGBColor(0x1E, 0x29, 0x3B))

    code_bg = add_shape(slide, x + Inches(0.3), Inches(2.7), Inches(1.2), Inches(0.45), m["color"])
    add_text(slide, x + Inches(0.3), Inches(2.73), Inches(1.2), Inches(0.35), m["code"],
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(3.5), Inches(0.35), m["name"],
             size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.65), Inches(3.5), Inches(0.35), m["subtitle"],
             size=10, color=m["color"])

    add_text(slide, x + Inches(0.3), Inches(4.1), Inches(3.5), Inches(2.9), m["desc"],
             size=11, color=GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10d: 멀티 에이전트 병렬 추론
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.5), "PARALLEL INFERENCE ENGINE", size=12, color=RGBColor(0xA8, 0x55, 0xF7), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "질문 1개를 여러 GPU가 협력 처리", size=32, color=WHITE, bold=True)

# 4가지 모드
modes = [
    {"name": "Single", "icon": "1️⃣",
     "desc": "간단한 질문\n에이전트 1개 처리\n최저 비용",
     "color": GRAY, "when": "질문 < 100자"},
    {"name": "Speculative", "icon": "🏎️",
     "desc": "N개 동시 시도\n가장 빠른 것 채택\n나머지 취소",
     "color": RGBColor(0xF5, 0x9E, 0x0B), "when": "보통 질문"},
    {"name": "Chunked", "icon": "✂️",
     "desc": "답변을 분할 병렬\n설명/코드/테스트\n속도 3배",
     "color": PRIMARY, "when": "복잡한 코딩"},
    {"name": "Tiered", "icon": "🔼",
     "desc": "7B가 초안 (3초)\n32B가 검증 (10초)\n품질=32B 속도=7B",
     "color": GREEN, "when": "고품질 필요"},
]

for i, m in enumerate(modes):
    x = Inches(0.5 + i * 3.2)
    card = add_shape(slide, x, Inches(2.3), Inches(3.0), Inches(2.5), RGBColor(0x1E, 0x29, 0x3B))

    add_text(slide, x + Inches(0.2), Inches(2.4), Inches(0.5), Inches(0.4), m["icon"], size=20)
    add_text(slide, x + Inches(0.8), Inches(2.4), Inches(2.0), Inches(0.35), m["name"],
             size=14, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.9), Inches(2.7), Inches(1.2), m["desc"],
             size=11, color=GRAY)
    add_text(slide, x + Inches(0.2), Inches(4.2), Inches(2.7), Inches(0.4),
             f"조건: {m['when']}", size=9, color=m["color"])

# 하단: 장애 처리 + 용량 적응
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.8), [
    {"text": "자동 장애 처리 + 용량 적응", "size": 14, "color": WHITE, "bold": True, "space": 10},
    {"text": "에이전트 장애 → 해당 청크만 다른 에이전트로 즉시 재할당 (응답 지연 최소화)", "size": 11, "color": GRAY, "space": 6},
    {"text": "RTX 4060 (8GB) → 7B 모델로 draft 담당  |  RTX 5090 (32GB) → V3로 검증 담당  |  자동 역할 분배", "size": 11, "color": CYAN, "space": 6},
    {"text": "모드 자동 선택: 질문 복잡도 + 에이전트 GPU tier + 현재 부하 → 최적 모드 결정", "size": 11, "color": RGBColor(0xF5, 0x9E, 0x0B), "space": 6},
    {"text": "하이브리드 (AI CDN): 평상시=서버 직접 처리 | 피크=에이전트 오버플로우 분배 | 점검=에이전트 전환 (무중단)", "size": 11, "color": GREEN, "space": 6},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 10e: HWARANG 토큰 이코노미
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.5), "HWARANG TOKEN ECONOMY", size=12, color=RGBColor(0xFB, 0xBF, 0x24), bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8), "3중 균형 토큰 이코노미", size=32, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
    "고정 상한 + 적응형 발행 + 자동 소각. 조작 불가능한 스마트 컨트랙트 기반.",
    size=14, color=GRAY)

# 3개 메커니즘 카드
mechanisms = [
    {
        "icon": "🔒", "name": "고정 상한",
        "desc": "총 10억 HWR\n절대 초과 불가\n스마트 컨트랙트에 하드코딩",
        "color": RGBColor(0xEF, 0x44, 0x44),
    },
    {
        "icon": "⚖️", "name": "적응형 발행",
        "desc": "GPU 부족 → 보상 ↑\nGPU 과잉 → 보상 ↓\n+ 성과 기반 반감기",
        "color": PRIMARY,
    },
    {
        "icon": "🔥", "name": "자동 소각",
        "desc": "AI 이용 시 30% 소각\n서비스 ↑ → 소각 ↑\n→ 디플레이션 → 가치 ↑",
        "color": RGBColor(0xF5, 0x9E, 0x0B),
    },
]

for i, m in enumerate(mechanisms):
    x = Inches(0.5 + i * 4.2)
    card = add_shape(slide, x, Inches(2.5), Inches(3.9), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, x + Inches(0.3), Inches(2.6), Inches(0.5), Inches(0.4), m["icon"], size=24)
    add_text(slide, x + Inches(1.0), Inches(2.6), Inches(2.6), Inches(0.4), m["name"],
             size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.2), Inches(3.4), Inches(1.3), m["desc"],
             size=12, color=GRAY)

# 하단: 시뮬레이션 결과
sim_box = add_shape(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.2), RGBColor(0x1E, 0x29, 0x3B))
add_multi_text(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(2.0), [
    {"text": "5년 시뮬레이션 결과", "size": 14, "color": WHITE, "bold": True, "space": 10},
    {"text": "Y1 (1천명): 일 발행 43만 HWR, 균형  |  Y3 (5만명): 디플레이션 시작  |  Y5 (50만명): 유통 5.9억 (40% 소각)", "size": 11, "color": GRAY, "space": 8},
    {"text": "GPU 보상: RTX 5090 기준 시간당 10 HWR × 적응형 계수 (GPU 부족 시 최대 3배)  |  학습 참여 시 2배", "size": 11, "color": CYAN, "space": 8},
    {"text": "조작 방지: 스마트 컨트랙트 → 발행/소각 규칙 변경 불가  |  DAO 거버넌스 (파라미터만 투표)", "size": 11, "color": RGBColor(0xF5, 0x9E, 0x0B), "space": 8},
])


# ════════════════════════════════════════════════════════════════
# SLIDE 11: 마무리
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_shape(slide, Inches(-2), Inches(-2), Inches(6), Inches(6), RGBColor(0x1E, 0x1B, 0x4B))
add_shape(slide, Inches(9), Inches(4), Inches(6), Inches(6), RGBColor(0x1E, 0x1B, 0x4B))

logo_shape = add_shape(slide, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5), PRIMARY)
add_text(slide, Inches(5.9), Inches(1.25), Inches(1.5), Inches(1.5), "H", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.0), Inches(9.3), Inches(0.8), "함께 만들어가는\n한국의 AI", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5), "나 혼자선 무리다. 같이 만들어가서, 같이 자산으로.", size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_multi_text(slide, Inches(3), Inches(5.5), Inches(7.3), Inches(1.5), [
    {"text": "hwarang.ai  |  admin.hwarang.ai", "size": 14, "color": CYAN, "align": PP_ALIGN.CENTER, "space": 8},
    {"text": "(주)퍼시스모어  |  hello@persismore.com", "size": 12, "color": GRAY, "align": PP_ALIGN.CENTER},
])


# ════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "..", "..", "Hwarang_AI_소개.pptx")
output_path = os.path.abspath(output_path)
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
